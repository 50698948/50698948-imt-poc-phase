"""
IMT PoC — FastAPI Backend Service

Start:  uvicorn api_server:app --reload --port 8000
"""

import json
import os
import sys
import time
import uuid
from contextlib import asynccontextmanager
from typing import Optional

from fastapi import FastAPI, HTTPException, Query, File, UploadFile, Form
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel

# Add parent poc/ directory to path so we can import existing modules
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "poc"))

from poc_standalone import (
    ingest,
    ingest_tickets_batch,
    get_ticket,
    get_ticket_by_id,
    update_ticket,
    retrieve,
    rerank,
    generate_action_plan,
    get_report_history_standalone,
    get_latest_report_standalone,
    get_recommendations_standalone,
    revise_task_standalone,
    _build_highlights,
    _build_report,
    _generate_recommendations,
    _save_recommendations,
    SEED_TICKETS,
    DB_PATH,
)


# ── Lifespan: clear DB on startup for clean demo state ──
@asynccontextmanager
async def lifespan(app: FastAPI):
    if os.path.exists(DB_PATH):
        os.remove(DB_PATH)
    yield


app = FastAPI(title="IMT PoC API", version="1.0", lifespan=lifespan)
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])

# ═══════════════════════════════════════════════════════════
# Request / Response Models
# ═══════════════════════════════════════════════════════════

class SeedRequest(BaseModel):
    pass

class RetrieveRequest(BaseModel):
    title: str
    description: str
    service_name: str = ""
    category: str = "application"
    severity: str = "P1"
    error_type: Optional[str] = None

class LifecycleRequest(BaseModel):
    pass

class ReviseTaskRequest(BaseModel):
    status: Optional[str] = None
    description: Optional[str] = None
    revision_note: Optional[str] = None
    revised_by: Optional[str] = None

# ═══════════════════════════════════════════════════════════
# Routes
# ═══════════════════════════════════════════════════════════

@app.get("/api/health")
def health():
    return {"status": "ok"}


@app.post("/api/seed")
def seed_database():
    """Seed the database with 49 mock tickets."""
    try:
        ids = ingest_tickets_batch(SEED_TICKETS)
        return {"status": "ok", "count": len(ids), "ids": ids}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/tickets")
def list_tickets(
    status: Optional[str] = None,
    category: Optional[str] = None,
    limit: int = Query(default=50, le=100),
):
    """List all tickets, optionally filtered."""
    import sqlite3
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    q = "SELECT id, incident_no, title, description, severity, service_name, category, status, error_type, version, created_at, updated_at FROM incident_tickets WHERE 1=1"
    params = []
    if status:
        q += " AND status=?"
        params.append(status)
    if category:
        q += " AND category=?"
        params.append(category)
    q += " ORDER BY created_at DESC LIMIT ?"
    params.append(limit)
    rows = conn.execute(q, params).fetchall()
    conn.close()
    return [dict(r) for r in rows]


@app.get("/api/tickets/{incident_no}")
def ticket_detail(incident_no: str):
    """Get a single ticket by incident_no."""
    t = get_ticket(incident_no)
    if t is None:
        raise HTTPException(status_code=404, detail="Ticket not found")
    return t


@app.post("/api/retrieve")
def run_retrieval(req: RetrieveRequest):
    """Run the full E2E retrieval pipeline and return results."""
    current = {
        "title": req.title,
        "description": req.description,
        "service_name": req.service_name,
        "category": req.category,
        "severity": req.severity,
        "error_type": req.error_type,
    }
    try:
        candidates = retrieve(current)
        reranked = rerank(current, candidates)
        plan = generate_action_plan(current, reranked)
        return {
            "current": current,
            "candidates": candidates[:15],
            "reranked": reranked[:5],
            "action_plan": plan,
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/lifecycle")
def run_lifecycle():
    """Run the lifecycle demo and return stage-by-stage results."""
    lifecycle_stages = [
        {"stage": "T+0min — Initial Alert (vague)", "incident_no": "INC-2025-0001",
         "title": "Payment service failing under load",
         "description": "Payment service returning errors. Users cannot complete checkout.",
         "severity": "P0", "service_name": "payment-service", "category": "application", "status": "open"},
        {"stage": "T+10min — Triage Complete", "incident_no": "INC-2025-0001", "status": "investigating",
         "description": "Payment service returning HTTP 504 Gateway Timeout on /charge endpoint. P99 latency jumped from 200ms to 30s. 60% of transactions failing. Thread dump shows 200 threads blocked on HikariCP getConnection(). Connection pool max=50, active=50, pending=340.",
         "error_type": "timeout"},
        {"stage": "T+45min — Root Cause Found", "incident_no": "INC-2025-0001", "status": "investigating",
         "description": "Payment service returning HTTP 504 Gateway Timeout on /charge endpoint. P99 latency 30s. Connection pool max=50, active=50, pending=340. Traced to a deployment that added a new fraud-check HTTP call inside the /charge handler. Each charge now makes 4 sequential HTTP calls, each holding a DB connection. Connection hold time doubled to 4s, exhausting the pool.",
         "error_type": "timeout",
         "root_cause": "New deployment added fraud-check HTTP call inside /charge handler, increasing critical path from 2 to 4 sequential calls. Each call holds a DB connection. Connection hold time doubled to ~4s, exhausting the 50-connection pool at 12 req/s."},
        {"stage": "T+90min — Resolved", "incident_no": "INC-2025-0001", "status": "resolved",
         "description": "Payment service returning HTTP 504 Gateway Timeout on /charge endpoint. P99 latency 30s. Connection pool exhausted at 50/50 with 340 pending.",
         "error_type": "timeout",
         "root_cause": "New deployment added fraud-check HTTP call inside /charge handler, increasing critical path to 4 sequential calls, doubling connection hold time.",
         "resolution": "1. Rolled back deployment. 2. Connection pool recovered within 30s. 3. Moved fraud-check to async post-processing. 4. Increased pool size to 150. 5. Added alert: HikariCP_pending_connections > 20."},
    ]

    results = []
    conn = __import__('sqlite3').connect(DB_PATH)

    for idx, stage in enumerate(lifecycle_stages):
        if idx == 0:
            ingest(stage)
            ticket = get_ticket(stage["incident_no"])
        else:
            ticket = update_ticket(stage["incident_no"],
                                   status=stage["status"],
                                   description=stage.get("description"),
                                   root_cause=stage.get("root_cause"),
                                   resolution=stage.get("resolution"),
                                   error_type=stage.get("error_type"))

        candidates = retrieve(ticket)
        candidates = [c for c in candidates if c.get("incident_no") != stage["incident_no"]]
        reranked = rerank(ticket, candidates)

        report = get_latest_report_standalone(stage["incident_no"])
        tasks = get_recommendations_standalone(stage["incident_no"])

        scores = [r.get("rerank_score", 0) for r in reranked[:5]]
        avg = sum(scores) / max(len(scores), 1) if scores else 0

        results.append({
            "stage": stage["stage"],
            "version": ticket.get("version", 1),
            "status": ticket.get("status", ""),
            "description": ticket.get("description", ""),
            "root_cause": ticket.get("root_cause", ""),
            "resolution": ticket.get("resolution", ""),
            "avg_rerank_score": round(avg, 1),
            "top5_reranked": [{
                "incident_no": r.get("incident_no"),
                "title": r.get("title"),
                "score": r.get("rerank_score"),
                "reason": r.get("rerank_reason"),
            } for r in reranked[:5]],
            "report_highlights": report.get("highlights", []) if report else [],
            "tasks": [{"task_order": t["task_order"], "description": t["description"],
                        "source": t.get("source"), "status": t.get("status"), "id": t["id"]}
                       for t in tasks],
        })

    return {"stages": results}


@app.get("/api/reports/{incident_no}")
def get_reports(incident_no: str, latest_only: bool = False):
    """Get report history or latest report."""
    if latest_only:
        r = get_latest_report_standalone(incident_no)
        if r is None:
            raise HTTPException(status_code=404, detail="No reports found")
        return r
    return get_report_history_standalone(incident_no)


@app.post("/api/reports/{incident_no}/generate")
def generate_report_for_ticket(incident_no: str):
    """Trigger report + task generation for a specific incident."""
    ticket = get_ticket(incident_no)
    if ticket is None:
        raise HTTPException(status_code=404, detail="Incident not found")

    from poc_standalone import _build_highlights, _build_report, _generate_recommendations, _save_recommendations
    import sqlite3, uuid, json

    candidates = retrieve(ticket)
    candidates = [c for c in candidates if c.get("incident_no") != incident_no]
    reranked = rerank(ticket, candidates)
    highlights = _build_highlights(ticket, reranked)
    content = _build_report(ticket, reranked, highlights)
    tasks = _generate_recommendations(ticket, reranked)

    conn = sqlite3.connect(DB_PATH)
    conn.execute(
        "INSERT INTO leader_reports(id, incident_no, ticket_version, content, highlights) VALUES(?,?,?,?,?)",
        (str(uuid.uuid4()), incident_no, ticket.get("version", 1), content, json.dumps(highlights)))
    _save_recommendations(conn, incident_no, ticket.get("version", 1), tasks)
    conn.commit()
    conn.close()

    return {"status": "ok", "incident_no": incident_no, "highlights": highlights,
            "tasks_count": len(tasks)}


class CreateIncidentRequest(BaseModel):
    title: str
    description: str
    severity: str
    service_name: str
    category: str = "application"
    error_type: Optional[str] = None
    source: str = "manual"


@app.post("/api/incidents/create")
def create_incident(req: CreateIncidentRequest):
    """Create a new incident ticket with automatic retrieval."""
    incident_no = f"INC-{time.strftime('%Y')}-{str(uuid.uuid4().int % 10000).zfill(4)}"
    ticket_data = {
        "incident_no": incident_no,
        "title": req.title,
        "description": req.description,
        "severity": req.severity,
        "service_name": req.service_name,
        "category": req.category,
        "status": "open",
        "error_type": req.error_type,
    }
    ingest(ticket_data)
    ticket = get_ticket(incident_no)
    if ticket is None:
        raise HTTPException(status_code=500, detail="Failed to create ticket")

    candidates = retrieve(ticket)
    reranked = rerank(ticket, candidates)

    return {
        "incident_no": incident_no,
        "ticket": ticket,
        "top5_similar": [{"incident_no": r.get("incident_no"), "title": r.get("title"), "score": r.get("rerank_score")} for r in reranked[:5]],
    }


class XMatterAlert(BaseModel):
    alertId: Optional[str] = None
    title: str
    description: str
    severity: str  # critical/warning/info
    affectedService: Optional[str] = None
    tags: Optional[list[str]] = None


@app.post("/api/alerts/xmatter")
def receive_xmatter_alert(alert: XMatterAlert):
    """Receive xMatter alert webhook — auto-create incident + retrieve similar."""
    sev_map = {"critical": "P0", "warning": "P1", "info": "P2"}
    severity = sev_map.get(alert.severity.lower(), "P1")
    service = alert.affectedService or "unknown"
    error_type = None
    if alert.tags:
        known = {"timeout", "OOM", "deadlock", "race_condition", "config_error", "dependency_failure", "resource_exhaustion", "auth_error", "rate_limit"}
        for t in alert.tags:
            if t in known: error_type = t; break

    incident_no = f"INC-{time.strftime('%Y')}-{str(uuid.uuid4().int % 10000).zfill(4)}"
    ingest({
        "incident_no": incident_no,
        "title": alert.title,
        "description": f"[xMatter Alert {alert.alertId or 'N/A'}] {alert.description}",
        "severity": severity,
        "service_name": service,
        "category": "application",
        "status": "open",
        "error_type": error_type,
    })
    ticket = get_ticket(incident_no)
    candidates = retrieve(ticket)
    reranked = rerank(ticket, candidates)

    return {
        "incident_no": incident_no,
        "source": "xmatter",
        "alert_id": alert.alertId,
        "top5_similar": [{"incident_no": r.get("incident_no"), "title": r.get("title"), "score": r.get("rerank_score")} for r in reranked[:5]],
        "summary": f"Created {incident_no}. Found {len(reranked)} similar cases."
    }


@app.get("/api/tasks/{incident_no}")
def get_tasks(incident_no: str):
    """Get recommended tasks for a ticket."""
    tasks = get_recommendations_standalone(incident_no)
    return tasks


@app.post("/api/chat/upload")
async def chat_upload(file: UploadFile, incident_no: str = Form(...)):
    """Upload a file (docx/pptx/txt) to parse and append to incident description."""
    content = await file.read()
    ext = file.filename.split(".")[-1].lower() if file.filename else "txt"
    parsed = ""

    try:
        if ext == "txt" or ext == "log":
            parsed = content.decode("utf-8", errors="replace")
        elif ext == "docx":
            from docx import Document as DocxDocument
            from io import BytesIO
            doc = DocxDocument(BytesIO(content))
            parsed = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif ext == "pptx":
            from pptx import Presentation
            from io import BytesIO
            prs = Presentation(BytesIO(content))
            parsed = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip())
        else:
            parsed = content.decode("utf-8", errors="replace")[:2000]
    except Exception as e:
        return {"error": f"Failed to parse {ext}: {str(e)}"}

    if not parsed.strip():
        return {"error": "No text content extracted"}

    ticket = get_ticket(incident_no)
    if ticket is None:
        return {"error": "Incident not found"}

    new_desc = f"{ticket.get('description','')}\n\n[via file upload: {file.filename}, {parsed[:1500]}]"
    update_ticket(incident_no, description=new_desc[:4000])
    updated = get_ticket(incident_no)

    candidates = retrieve(updated)
    reranked = rerank(updated, [c for c in candidates if c.get("incident_no") != incident_no])

    old_len = len(ticket.get("description", ""))
    new_len = len(updated.get("description", ""))

    return {
        "parsed_text": parsed[:500],
        "word_count": len(parsed.split()),
        "desc_before": old_len,
        "desc_after": new_len,
        "filename": file.filename,
        "top_match": reranked[0].get("incident_no") if reranked else "N/A",
    }


@app.post("/api/reports/{incident_no}/publish")


@app.post("/api/tasks/{task_id}/revise")
def revise_task(task_id: str, req: ReviseTaskRequest):
    """Revise a recommended task with history recording."""
    fields = {k: v for k, v in req.model_dump().items() if v is not None}
    if not fields:
        raise HTTPException(status_code=400, detail="No fields to revise")
    result = revise_task_standalone(task_id, **fields)
    if result is None:
        raise HTTPException(status_code=404, detail="Task not found")

    # Record revision history
    import sqlite3, uuid as _uuid
    action = fields.get("description") and "modified" or fields.get("status", "updated")
    detail = fields.get("revision_note", "") or fields.get("description", "") or fields.get("status", "")
    conn = sqlite3.connect(DB_PATH)
    conn.execute(
        "INSERT INTO task_revision_history(id, task_id, incident_no, action, revised_by, detail) VALUES(?,?,?,?,?,?)",
        (str(_uuid.uuid4()), task_id, result.get("incident_no", ""), action,
         fields.get("revised_by", ""), str(detail)[:200]))
    conn.commit()
    conn.close()
    return result


@app.get("/api/tasks/{incident_no}/history")
def get_task_history(incident_no: str):
    """Get revision history for all tasks of an incident."""
    import sqlite3
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    rows = conn.execute(
        "SELECT * FROM task_revision_history WHERE incident_no=? ORDER BY created_at DESC LIMIT 30",
        (incident_no,)).fetchall()
    conn.close()
    return [{"id": r["id"], "task_id": r["task_id"], "action": r["action"],
             "revised_by": r["revised_by"], "detail": r["detail"],
             "created_at": r["created_at"]} for r in rows]


class ChatMessage(BaseModel):
    message: str
    incident_no: str = ""


@app.post("/api/chat")
def chat_endpoint(req: ChatMessage):
    """Chat interface — natural language incident management."""
    msg = req.message.strip().lower()
    incident_no = req.incident_no

    if not incident_no:
        return {"reply": "Please specify an incident number (e.g., INC-2025-0001).", "actions": []}

    ticket = get_ticket(incident_no)
    if ticket is None:
        return {"reply": f"Incident {incident_no} not found.", "actions": []}

    actions = []
    status_map = {"open": "open", "investigating": "investigating", "mitigated": "mitigated", "resolved": "resolved"}
    for k, v in status_map.items():
        if k in msg:
            update_ticket(incident_no, status=v)
            ticket = get_ticket(incident_no)
            actions.append(f"status→{v}")
            break

    if "root cause" in msg:
        idx = msg.find("root cause")
        rc = msg[idx:].split("\n")[0]
        rc = rc.replace("root cause", "").strip().lstrip(": is ").strip()
        if len(rc) > 10:
            update_ticket(incident_no, root_cause=rc[:500])
            ticket = get_ticket(incident_no)
            actions.append("root_cause:updated")

    if "resolution" in msg or "resolved by" in msg:
        for kw in ["resolution:", "resolved by:", "fix:"]:
            if kw in msg:
                idx = msg.find(kw)
                res = msg[idx + len(kw):].split("\n")[0].strip()
                if len(res) > 10:
                    update_ticket(incident_no, resolution=res[:500])
                    ticket = get_ticket(incident_no)
                    actions.append("resolution:updated")
                break

    if "description" in msg and (":" in msg or "update" in msg):
        idx = msg.find("description")
        desc = msg[idx:].split("\n")[0]
        desc = desc.replace("description", "").strip().lstrip(": is ").strip()
        if len(desc) > 10:
            update_ticket(incident_no, description=desc[:1000])
            ticket = get_ticket(incident_no)
            actions.append("description:updated")

    # Build reply
    reply = [
        f"**{ticket.get('incident_no')}** — {ticket.get('title','')}",
        f"Status: **{ticket.get('status','open').upper()}** | Severity: {ticket.get('severity','?')} | v{ticket.get('version',1)}",
    ]
    if ticket.get("root_cause"):
        reply.append(f"\nRoot Cause: {ticket['root_cause'][:200]}")
    if ticket.get("resolution"):
        reply.append(f"Resolution: {ticket['resolution'][:200]}")
    if actions:
        reply.append(f"\n✅ Actions: {', '.join(actions)}")

    if "recommend" in msg or "task" in msg:
        tasks = get_recommendations_standalone(incident_no)
        if tasks:
            reply.append(f"\n**Recommendations ({len(tasks)}):**")
            icons = {"pending": "○", "in_progress": "◐", "completed": "●", "rejected": "✕"}
            for t in tasks[:5]:
                reply.append(f"  {icons.get(t['status'],'?')} T{t['task_order']:02d} {t['description'][:80]}")

    if "report" in msg or "summary" in msg:
        report = get_latest_report_standalone(incident_no)
        if report:
            reply.append("\n**Latest Report:**")
            for hl in report.get("highlights", [])[:3]:
                reply.append(f"  * {hl}")

    if not actions and not any(kw in msg for kw in ["recommend", "task", "report", "summary"]):
        reply.append("\n💡 Try: 'update status to investigating', 'root cause: ...', 'recommendations', 'latest report'")

    return {"reply": "\n".join(reply), "actions": actions}


@app.get("/api/incidents/{incident_no}/timeline")
def incident_timeline(incident_no: str):
    """Chronological timeline of incident events."""
    ticket = get_ticket(incident_no)
    if ticket is None:
        raise HTTPException(status_code=404, detail="Incident not found")

    events = [{
        "time": ticket.get("created_at", ""), "type": "created", "icon": "○",
        "title": "Ticket Created", "detail": f"Severity: {ticket.get('severity','?')} | {ticket.get('service_name','')}"
    }]

    reports = get_report_history_standalone(incident_no)
    for r in reports:
        events.append({
            "time": r.get("generated_at", ""), "type": "report", "icon": "📄",
            "title": f"Report Demo v{r['ticket_version']}",
            "detail": " | ".join(r.get("highlights", [])[:2])
        })

    seen = set()
    for r in reports:
        v = r["ticket_version"]
        if v not in seen:
            seen.add(v)
            tasks = get_recommendations_standalone(incident_no, ticket_version=v)
            if tasks:
                events.append({
                    "time": "", "type": "tasks", "icon": "📋",
                    "title": f"{len(tasks)} Tasks Generated",
                    "detail": f"Version v{v}"
                })

    events.sort(key=lambda e: e["time"] or "", reverse=True)
    return {"incident_no": incident_no, "title": ticket.get("title", ""),
            "status": ticket.get("status", ""), "version": ticket.get("version", 1), "events": events}


class PublishReportRequest(BaseModel):
    revised_highlights: Optional[list[str]] = None
    executive_summary: Optional[str] = None
    review_notes: Optional[str] = None
    published_by: str = "engineer"


@app.post("/api/reports/{incident_no}/publish")
def publish_report(incident_no: str, req: PublishReportRequest):
    """Publish a draft report with optional revisions."""
    import sqlite3
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    row = conn.execute(
        "SELECT * FROM leader_reports WHERE incident_no=? ORDER BY ticket_version DESC LIMIT 1",
        (incident_no,)).fetchone()
    if not row:
        conn.close()
        raise HTTPException(status_code=404, detail="No report found")

    now = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())
    updates = {"report_status": "published", "published_at": now, "published_by": req.published_by}
    if req.review_notes:
        updates["review_notes"] = req.review_notes
        updates["reviewed_by"] = req.published_by
        updates["reviewed_at"] = now

    if req.revised_highlights is not None:
        import json
        updates["highlights"] = json.dumps(req.revised_highlights)
        existing = json.loads(row["revised_fields"] or "[]")
        existing.append("highlights")
        updates["revised_fields"] = json.dumps(existing)

    if req.executive_summary is not None:
        content = row["content"]
        parts = content.split("## 2. Current Status")
        if len(parts) > 1:
            content = f"## 1. Executive Summary\n{req.executive_summary}\n\n## 2. Current Status{parts[1]}"
        updates["content"] = content
        existing = json.loads(row["revised_fields"] or "[]")
        existing.append("executive_summary")
        updates["revised_fields"] = json.dumps(existing)

    set_clause = ", ".join(f"{k}=?" for k in updates)
    conn.execute(f"UPDATE leader_reports SET {set_clause} WHERE id=?", list(updates.values()) + [row["id"]])
    conn.commit()
    conn.close()
    return {"status": "ok", "report_status": "published"}


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
